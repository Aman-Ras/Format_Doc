from flask import Flask, request, render_template_string, send_file, flash, jsonify
from dotenv import load_dotenv
import io
import os
import re
import json
import uuid
import zipfile
import threading
import requests
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT as WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
import PyPDF2
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
try:
    from pyngrok import ngrok
except Exception:
    ngrok = None

load_dotenv()
api_key = os.getenv('OPENAI_API_KEY')

app = Flask(__name__)
app.secret_key = 'groq_resume_formatter_2025'

# In-memory job store: job_id -> {total, done, errors, zip_buffer, zip_name, lock}
jobs = {}
jobs_lock = threading.Lock()

# Network configuration for local network access
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Clean HTML Template
UPLOAD_FORM = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Resume Formatter — PowerPoint Generator</title>
    <meta name="description" content="Upload your resume and get a professionally formatted PowerPoint presentation instantly.">
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

        body {
            font-family: 'Inter', sans-serif;
            background: #f0f2f5;
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 24px;
        }

        .page-wrapper {
            width: 100%;
            max-width: 520px;
        }

        /* ── Brand mark ── */
        .brand {
            text-align: center;
            margin-bottom: 28px;
        }
        .brand-icon {
            width: 52px;
            height: 52px;
            background: linear-gradient(135deg, #1a73e8 0%, #0d47a1 100%);
            border-radius: 14px;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            margin-bottom: 14px;
            box-shadow: 0 4px 14px rgba(26,115,232,0.35);
        }
        .brand-icon svg { width: 26px; height: 26px; fill: #fff; }
        .brand h1 {
            font-size: 22px;
            font-weight: 700;
            color: #111827;
            letter-spacing: -0.4px;
        }
        .brand p {
            font-size: 13.5px;
            color: #6b7280;
            margin-top: 4px;
            font-weight: 400;
        }

        /* ── Card ── */
        .card {
            background: #ffffff;
            border-radius: 20px;
            padding: 36px 36px 32px;
            box-shadow: 0 1px 3px rgba(0,0,0,0.06), 0 8px 32px rgba(0,0,0,0.08);
            border: 1px solid rgba(0,0,0,0.05);
        }

        /* ── Error banner ── */
        .error-banner {
            background: #fff5f5;
            border: 1px solid #fed7d7;
            border-radius: 10px;
            padding: 12px 16px;
            margin-bottom: 20px;
            color: #c53030;
            font-size: 13.5px;
            font-weight: 500;
        }

        /* ── Upload Zone ── */
        .upload-zone {
            border: 2px dashed #d1d5db;
            border-radius: 14px;
            padding: 28px 20px;
            text-align: center;
            cursor: pointer;
            transition: border-color 0.2s, background 0.2s;
            position: relative;
            background: #fafafa;
        }
        .upload-zone:hover, .upload-zone.drag-over {
            border-color: #1a73e8;
            background: #f0f6ff;
        }
        .upload-zone input[type="file"] {
            position: absolute;
            inset: 0;
            opacity: 0;
            cursor: pointer;
            width: 100%;
            height: 100%;
        }
        .upload-icon {
            width: 44px;
            height: 44px;
            background: #e8f0fe;
            border-radius: 50%;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            margin-bottom: 12px;
        }
        .upload-icon svg { width: 22px; height: 22px; fill: #1a73e8; }
        .upload-label {
            font-size: 14px;
            font-weight: 600;
            color: #111827;
            margin-bottom: 4px;
        }
        .upload-hint {
            font-size: 12.5px;
            color: #9ca3af;
        }
        .upload-hint span { color: #1a73e8; font-weight: 500; }

        /* ── File chosen state ── */
        .file-chosen {
            display: none;
            align-items: center;
            gap: 10px;
            background: #f0f6ff;
            border: 1.5px solid #93c5fd;
            border-radius: 10px;
            padding: 10px 14px;
            margin-top: 14px;
        }
        .file-chosen.visible { display: flex; }
        .file-chosen-icon { font-size: 20px; }
        .file-chosen-name {
            font-size: 13px;
            font-weight: 500;
            color: #1e3a6e;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
            flex: 1;
        }
        .file-chosen-remove {
            cursor: pointer;
            color: #6b7280;
            font-size: 18px;
            line-height: 1;
            flex-shrink: 0;
        }
        .file-chosen-remove:hover { color: #ef4444; }

        /* ── Supported formats ── */
        .formats {
            display: flex;
            gap: 6px;
            flex-wrap: wrap;
            margin-top: 16px;
            justify-content: center;
        }
        .format-tag {
            font-size: 11px;
            font-weight: 500;
            color: #6b7280;
            background: #f3f4f6;
            border-radius: 20px;
            padding: 3px 10px;
            letter-spacing: 0.3px;
        }

        /* ── Divider ── */
        .divider { height: 1px; background: #f0f0f0; margin: 24px 0; }

        /* ── Submit button ── */
        .submit-btn {
            width: 100%;
            background: linear-gradient(135deg, #1a73e8 0%, #1557b0 100%);
            color: #fff;
            border: none;
            border-radius: 12px;
            padding: 14px 24px;
            font-size: 15px;
            font-weight: 600;
            font-family: 'Inter', sans-serif;
            cursor: pointer;
            letter-spacing: 0.2px;
            transition: transform 0.15s, box-shadow 0.15s, opacity 0.2s;
            box-shadow: 0 4px 14px rgba(26,115,232,0.35);
            display: flex;
            align-items: center;
            justify-content: center;
            gap: 8px;
        }
        .submit-btn:hover:not(:disabled) {
            transform: translateY(-1px);
            box-shadow: 0 6px 20px rgba(26,115,232,0.45);
        }
        .submit-btn:active:not(:disabled) { transform: translateY(0); }
        .submit-btn:disabled { opacity: 0.65; cursor: not-allowed; transform: none; }
        .btn-icon { width: 18px; height: 18px; fill: #fff; flex-shrink: 0; }

        /* ── Footer note ── */
        .footer-note {
            text-align: center;
            margin-top: 20px;
            font-size: 12px;
            color: #9ca3af;
        }

        /* ── Loading overlay ── */
        .loading-overlay {
            position: fixed;
            inset: 0;
            background: rgba(255,255,255,0.88);
            backdrop-filter: blur(4px);
            display: none;
            align-items: center;
            justify-content: center;
            z-index: 9999;
        }
        .loading-card {
            background: #fff;
            border-radius: 20px;
            padding: 36px 44px;
            text-align: center;
            box-shadow: 0 8px 40px rgba(0,0,0,0.12);
            border: 1px solid #e5e7eb;
        }
        .spinner-ring {
            width: 52px;
            height: 52px;
            border: 4px solid #e8f0fe;
            border-top-color: #1a73e8;
            border-radius: 50%;
            animation: spin 0.9s linear infinite;
            margin: 0 auto 16px;
        }
        @keyframes spin { to { transform: rotate(360deg); } }
        .loading-title {
            font-size: 15px;
            font-weight: 600;
            color: #111827;
            margin-bottom: 4px;
        }
        .loading-sub {
            font-size: 12.5px;
            color: #9ca3af;
        }
    </style>
</head>
<body>
    <div class="page-wrapper">
        <!-- Brand -->
        <div class="brand">
            <div class="brand-icon">
                <svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8l-6-6zm-1 1.5L18.5 9H13V3.5zM6 20V4h5v7h7v9H6z"/></svg>
            </div>
            <h1>Resume Formatter</h1>
            <p>Upload your resume — get a polished PowerPoint instantly</p>
        </div>

        <!-- Card -->
        <div class="card">

            {% with messages = get_flashed_messages() %}
                {% if messages %}
                    <div class="error-banner">
                        {% for message in messages %}⚠ {{ message }}{% endfor %}
                    </div>
                {% endif %}
            {% endwith %}

            <form id="uploadForm" method="POST" enctype="multipart/form-data">

                <!-- Upload zone -->
                <div class="upload-zone" id="uploadZone">
                    <input type="file" id="resume_file" name="resume_file"
                           accept=".pdf,.doc,.docx,.txt" required multiple>
                    <div class="upload-icon">
                        <svg viewBox="0 0 24 24"><path d="M19.35 10.04A7.49 7.49 0 0 0 12 4C9.11 4 6.6 5.64 5.35 8.04A5.994 5.994 0 0 0 0 14c0 3.31 2.69 6 6 6h13c2.76 0 5-2.24 5-5 0-2.64-2.05-4.78-4.65-4.96zM14 13v4h-4v-4H7l5-5 5 5h-3z"/></svg>
                    </div>
                    <p class="upload-label">Drop your resumes here</p>
                    <p class="upload-hint">or <span>browse to upload</span> · supports multiple files</p>
                </div>

                <!-- Selected file indicator -->
                <div class="file-chosen" id="fileChosen">
                    <span class="file-chosen-icon">📄</span>
                    <span class="file-chosen-name" id="fileChosenName"></span>
                    <span class="file-chosen-remove" id="fileRemove" title="Remove">✕</span>
                </div>

                <!-- Supported formats -->
                <div class="formats">
                    <span class="format-tag">PDF</span>
                    <span class="format-tag">DOCX</span>
                    <span class="format-tag">DOC</span>
                    <span class="format-tag">TXT</span>
                </div>

                <div class="divider"></div>

                <button id="submitBtn" type="submit" class="submit-btn">
                    <svg class="btn-icon" viewBox="0 0 24 24"><path d="M19 9h-4V3H9v6H5l7 7 7-7zM5 18v2h14v-2H5z"/></svg>
                    Generate PowerPoint Resume
                </button>

            </form>
        </div>

        <p class="footer-note">Your resume is processed securely and never stored.</p>
    </div>

    <!-- Loading overlay -->
    <div id="loadingOverlay" class="loading-overlay">
        <div class="loading-card">
            <div class="spinner-ring"></div>
            <p class="loading-title">Generating your resume…</p>
            <p class="loading-sub">This may take 10–20 seconds</p>
        </div>
    </div>

<script>
(function () {
    var form     = document.getElementById('uploadForm');
    var input    = document.getElementById('resume_file');
    var zone     = document.getElementById('uploadZone');
    var chosen   = document.getElementById('fileChosen');
    var chosenName = document.getElementById('fileChosenName');
    var removeBtn  = document.getElementById('fileRemove');
    var overlay  = document.getElementById('loadingOverlay');
    var btn      = document.getElementById('submitBtn');

    // File selection feedback
    input.addEventListener('change', function () {
        if (input.files && input.files.length > 0) {
            var count = input.files.length;
            chosenName.textContent = count === 1
                ? input.files[0].name
                : count + ' files selected';
            chosen.classList.add('visible');
        }
    });

    // Remove file
    removeBtn.addEventListener('click', function (e) {
        e.preventDefault();
        input.value = '';
        chosen.classList.remove('visible');
        chosenName.textContent = '';
    });

    // Drag-over highlight
    ['dragenter','dragover'].forEach(function(evt) {
        zone.addEventListener(evt, function(e) { e.preventDefault(); zone.classList.add('drag-over'); });
    });
    ['dragleave','drop'].forEach(function(evt) {
        zone.addEventListener(evt, function(e) { e.preventDefault(); zone.classList.remove('drag-over'); });
    });
    zone.addEventListener('drop', function(e) {
        var files = e.dataTransfer && e.dataTransfer.files;
        if (files && files.length > 0) {
            input.files = files;
            chosenName.textContent = files[0].name;
            chosen.classList.add('visible');
        }
    });

    // Form submit
    form.addEventListener('submit', function (e) {
        e.preventDefault();
        btn.disabled = true;
        btn.innerHTML = '<svg class="btn-icon" viewBox="0 0 24 24" style="fill:#fff"><path d="M12 4V1L8 5l4 4V6c3.31 0 6 2.69 6 6 0 1.01-.25 1.97-.7 2.8l1.46 1.46A7.93 7.93 0 0 0 20 12c0-4.42-3.58-8-8-8zm0 14c-3.31 0-6-2.69-6-6 0-1.01.25-1.97.7-2.8L5.24 7.74A7.93 7.93 0 0 0 4 12c0 4.42 3.58 8 8 8v3l4-4-4-4v3z"/></svg> Generating…';
        overlay.style.display = 'flex';

        var formData = new FormData(form);
        fetch('/', { method: 'POST', body: formData })
            .then(function (response) {
                var disposition  = response.headers.get('Content-Disposition') || '';
                var contentType  = (response.headers.get('Content-Type') || '').toLowerCase();
        var isPptx = disposition.indexOf('attachment') !== -1
                  || contentType.indexOf('presentationml') !== -1
                  || contentType.indexOf('zip') !== -1;

                if (isPptx) {
                    var filename = input.files.length > 1 ? 'formatted_resumes.zip' : 'formatted_resume.pptx';
                    var match = /filename\*=UTF-8''([^;]+)|filename="?([^";]+)"?/i.exec(disposition);
                    if (match) filename = decodeURIComponent(match[1] || match[2] || filename);
                    return response.blob().then(function (blob) {
                        var url = window.URL.createObjectURL(blob);
                        var a   = document.createElement('a');
                        a.href = url; a.download = filename;
                        document.body.appendChild(a); a.click(); a.remove();
                        window.URL.revokeObjectURL(url);
                        overlay.style.display = 'none';
                        btn.disabled = false;
                        btn.innerHTML = '<svg class="btn-icon" viewBox="0 0 24 24" style="fill:#fff"><path d="M19 9h-4V3H9v6H5l7 7 7-7zM5 18v2h14v-2H5z"/></svg> Generate PowerPoint Resume';
                    });
                }
                return response.text().then(function (html) {
                    document.open(); document.write(html); document.close();
                });
            })
            .catch(function () {
                overlay.style.display = 'none';
                btn.disabled = false;
                btn.innerHTML = '<svg class="btn-icon" viewBox="0 0 24 24" style="fill:#fff"><path d="M19 9h-4V3H9v6H5l7 7 7-7zM5 18v2h14v-2H5z"/></svg> Generate PowerPoint Resume';
                alert('Network error. Please try again.');
            });
    });
})();
</script>
</body>
</html>
'''


# Helper Functions for DOCX Formatting
def set_cell_background_white(cell):
    """Set table cell background to white"""
    shading = OxmlElement('w:shd')
    shading.set(qn('w:val'), 'clear')
    shading.set(qn('w:color'), 'auto')
    shading.set(qn('w:fill'), 'FFFFFF')
    cell._element.get_or_add_tcPr().append(shading)

def set_calibri_font(run, size=11, bold=False, color=RGBColor(0, 0, 0)):
    """Set Calibri font with consistent styling"""
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

def remove_table_borders(table):
    """Remove table borders for clean look"""
    try:
        for row in table.rows:
            for cell in row.cells:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                tcBorders = tcPr.first_child_found_in("w:tcBorders")
                if tcBorders is not None:
                    tcPr.remove(tcBorders)
    except:
        pass

# OpenAI API Resume Extractor Class
class OpenAIResumeExtractor:
    def __init__(self, api_key=None):
        self.api_key = api_key or os.getenv('OPENAI_API_KEY')
        self.api_url = "https://api.openai.com/v1/chat/completions"
        # Default to a fast JSON-capable model; adjust as needed
        self.model = "gpt-4.1-mini"
    
    def _clean_json_text(self, text):
        """Clean model text to maximize chances of valid JSON parsing."""
        if not isinstance(text, str):
            return "{}"
        cleaned = text.strip()
        # Strip fenced code blocks ```json ... ``` or ``` ... ```
        if cleaned.startswith("```"):
            cleaned = re.sub(r"^```(?:json)?\n", "", cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(r"\n```$", "", cleaned)
        # Extract content between the first '{' and the last '}' if present
        json_start = cleaned.find('{')
        json_end = cleaned.rfind('}')
        if json_start != -1 and json_end != -1 and json_end > json_start:
            cleaned = cleaned[json_start:json_end+1]
        # Normalize smart quotes
        cleaned = cleaned.replace('\u201c', '"').replace('\u201d', '"').replace('\u2019', "'")
        cleaned = cleaned.replace('“', '"').replace('”', '"').replace('’', "'")
        # Remove trailing commas before closing braces/brackets
        cleaned = re.sub(r",\s*(\}|\])", r"\1", cleaned)
        return cleaned
    
    def extract_with_openai(self, resume_text):
        """Extract structured data using OpenAI API"""
        if not self.api_key:
            print("OPENAI API key not provided, falling back to rule-based extraction")
            return self.extract_with_rules(resume_text)
        
        # Enhanced prompt for comprehensive extraction
        prompt = f"""
        You are a professional resume parser. Extract structured information from this resume text and return ONLY valid JSON.

        Resume Text:
        {resume_text}

        Extract and return a JSON object with these exact fields:
        {{
            "name": "Full candidate name",
            "location": "City/Location mentioned",
            "email": "Email address if found, otherwise empty string",
            "phone": "Phone number if found, otherwise empty string", 
            "date": "Application date if found, otherwise current date in DD-MMM-YYYY format",
            "subject": "Subject line in the exact format 'Application for the Position of {{ROLE}}' (no leading 'Subject:' label)",
            "summary": "Professional summary/objective from Key Expertise section and first paragraph. Write a COMPLETE summary that ends naturally at a full sentence — do NOT truncate or add '...'. Keep the total within 500 characters.",
            "education": [
                {{"degree": "Degree name", "institution": "Institution name"}}
            ],
            "experience_table": [
                {{
                    "company_name": "Company Name as Role (Duration)",
                    "roles_responsibility": [
                        "Exact responsibility 1 as written in resume",
                        "Exact responsibility 2 as written in resume",
                        "Exact responsibility 3 as written in resume"
                    ]
                }}
            ],
            "skills": ["List of technical skills, tools, technologies from Key Skills section"],
            "certifications": ["List of certifications and qualifications"],
            "cover_letter": "Complete cover letter content including all paragraphs after 'Dear Hiring Manager,'. If no cover letter text exists in the resume, GENERATE a professional, concise cover letter tailored to the 'subject' and the candidate's profile, using proper paragraph structure (2-5 paragraphs)."
        }}
        
        CRITICAL RULES:
        1. Extract ONLY information clearly present in the text
        2. For "name": Extract the candidate's full name as shown in the document header
        3. For "location": Extract city/location (e.g., "Bengaluru", "Bangalore")
        4. For "date": Extract application date if found, otherwise use current date
        5. For "subject": Return exactly "Application for the Position of {{ROLE}}" using the candidate's target role/title; do not include a leading "Subject:" label or extra punctuation
        6. For "summary": Combine Key Expertise section and first paragraph of cover letter into a concise professional summary within 500 characters. The summary MUST end at a complete sentence — never cut mid-word or add '...'
        7. For "education": Extract degree and institution in table format
        8. For "experience_table": Format each entry to match the table structure with "Company Name as Role (Duration)" in company_name field
        9. For "roles_responsibility" field: Extract each responsibility EXACTLY as written by the candidate from the Roles & Responsibility column
        10. For "skills": Extract from Key Skills section as a bulleted list
        11. For "certifications": Extract from CERTIFICATIONS section
        12. For "cover_letter": GENERATE a professional cover letter using the candidate's details (summary, key skills, certifications, experience highlights) tailored to the subject, and return it as multi-paragraph text (avoid repeating "Dear Hiring Manager,").
        13. Preserve the original formatting, punctuation, and exact language used by the candidate
        14. Return ONLY valid JSON with proper array formatting
        15. Ensure all JSON syntax is correct with proper quotes and commas
        """

        
        headers = {
            'Authorization': f'Bearer {self.api_key}',
            'Content-Type': 'application/json'
        }
        
        data = {
            'model': self.model,
            'messages': [
                {'role': 'system', 'content': 'You are a professional resume parsing expert. Extract information accurately and return only valid JSON.'},
                {'role': 'user', 'content': prompt}
            ],
            'temperature': 0.1,
            'max_tokens': 2500,
            'top_p': 1,
            'stream': False,
            # Enforce pure JSON response when supported
            'response_format': {'type': 'json_object'}
        }
        
        try:
            response = requests.post(self.api_url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            extracted_text = result['choices'][0]['message']['content'].strip()
            # Clean and parse JSON response
            json_str = self._clean_json_text(extracted_text)
            try:
                extracted_data = json.loads(json_str)
                return self.validate_and_clean_data(extracted_data)
            except json.JSONDecodeError as inner_e:
                print(f"JSON parsing error after cleaning: {str(inner_e)}")
                print("Raw model output (truncated):", extracted_text[:1000])
                return self.extract_with_rules(resume_text)
                
        except requests.exceptions.RequestException as e:
            print(f"OpenAI API request error: {str(e)}")
            return self.extract_with_rules(resume_text)
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {str(e)}")
            return self.extract_with_rules(resume_text)
        except Exception as e:
            print(f"OpenAI extraction error: {str(e)}")
            return self.extract_with_rules(resume_text)
    
    def validate_and_clean_data(self, data):
        """Validate and clean extracted data"""
        defaults = {
            'name': '',
            'location': '',
            'email': '',
            'phone': '',
            'date': '',
            'subject': '',
            'summary': '',
            'education': [],
            'experience_table': [],
            'experience': [],  # Keep for backward compatibility
            'skills': [],
            'certifications': [],
            'cover_letter': ''
        }
        
        # Convert education format if needed
        if 'education' in data and isinstance(data['education'], list):
            education_tuples = []
            for edu in data['education']:
                if isinstance(edu, dict):
                    degree = edu.get('degree', 'BCA')
                    institution = edu.get('institution', 'University Name')
                    education_tuples.append((degree, institution))
                elif isinstance(edu, (list, tuple)) and len(edu) >= 2:
                    education_tuples.append((edu[0], edu[1]))
            data['education'] = education_tuples if education_tuples else defaults['education']
        
        # Merge with defaults for missing fields
        for key, default_value in defaults.items():
            if key not in data or not data[key]:
                data[key] = default_value
        
        # Clean and limit arrays
        if isinstance(data.get('skills'), list):
            data['skills'] = [skill.strip() for skill in data['skills'][:15] if isinstance(skill, str) and skill.strip()]
        
        if isinstance(data.get('certifications'), list):
            data['certifications'] = [cert.strip() for cert in data['certifications'][:10] if isinstance(cert, str) and cert.strip()]

        
        # Use LLM-provided cover letter as-is (trim only)
        if isinstance(data.get('cover_letter'), str):
            data['cover_letter'] = data['cover_letter'].strip()
        else:
            data['cover_letter'] = ''
        
        # Ensure cover letter ends with desired closing
        if data['cover_letter']:
            candidate_name = (data.get('name') or '').strip() or 'Candidate'
            desired_closing = f"Sincerely,\n{candidate_name}"
            current = data['cover_letter'].rstrip()
            if not current.endswith(desired_closing):
                data['cover_letter'] = current + "\n\n" + desired_closing

        return data

    def extract_with_rules(self, resume_text):
        """Section-aware rule-based extraction"""
        lines = [line.strip() for line in resume_text.split('\n') if line.strip()]

        info = {
            'name': 'Professional Candidate',
            'location': '',
            'email': '',
            'phone': '',
            'date': datetime.now().strftime("%d-%b-%Y"),
            'subject': 'Application for Technical Position',
            'summary': '',
            'education': [],
            'experience_table': [],
            'experience': [],
            'skills': [],
            'certifications': []
        }

        # --- Name: first clean line with no digits/@ ---
        for line in lines[:5]:
            if (2 < len(line.split()) <= 5 and
                not any(c.isdigit() for c in line[:15]) and
                '@' not in line and '|' not in line):
                info['name'] = line
                break

        # --- Contact ---
        email_match = re.search(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b', resume_text)
        if email_match:
            info['email'] = email_match.group()

        phone_match = re.search(r'\+?\d[\d\s\-\(\)]{9,14}', resume_text)
        if phone_match:
            info['phone'] = phone_match.group().strip()

        location_cities = ['bengaluru', 'bangalore', 'mumbai', 'delhi', 'chennai',
                           'hyderabad', 'pune', 'india', 'jaipur', 'noida', 'kolkata']
        for line in lines[:10]:
            if any(c in line.lower() for c in location_cities) and len(line) < 120:
                info['location'] = line.split('|')[0].strip()
                break

        # --- Section splitting ---
        SECTION_HEADERS = {
            'summary':        ['PROFESSIONAL SUMMARY', 'SUMMARY', 'OBJECTIVE', 'PROFILE'],
            'experience':     ['WORK EXPERIENCE', 'EXPERIENCE', 'EMPLOYMENT HISTORY'],
            'education':      ['EDUCATION', 'ACADEMIC BACKGROUND'],
            'skills':         ['SKILLS', 'TECHNICAL SKILLS', 'KEY SKILLS', 'CORE COMPETENCIES'],
            'certifications': ['CERTIFICATIONS & ACHIEVEMENTS', 'CERTIFICATIONS', 'ACHIEVEMENTS', 'CERTIFICATION'],
            'projects':       ['PROJECTS', 'PERSONAL PROJECTS'],
        }

        # Build ordered list of (section_name, char_position)
        upper_text = resume_text.upper()
        found = []
        for sec, markers in SECTION_HEADERS.items():
            for m in markers:
                for pat in (f'\n{m}\n', f'\n{m} \n', f'\n\n{m}'):
                    pos = upper_text.find(pat)
                    if pos != -1:
                        found.append((sec, pos))
                        break
        found.sort(key=lambda x: x[1])
        # Deduplicate section names keeping first occurrence
        seen = set()
        ordered = []
        for sec, pos in found:
            if sec not in seen:
                seen.add(sec)
                ordered.append((sec, pos))

        def get_section(name):
            names = [s for s, _ in ordered]
            poses = [p for _, p in ordered]
            if name not in names:
                return ''
            i = names.index(name)
            start = poses[i]
            end = poses[i + 1] if i + 1 < len(poses) else len(resume_text)
            return resume_text[start:end]

        # --- Summary ---
        sec = get_section('summary')
        if sec:
            body = [l.strip() for l in sec.split('\n')
                    if l.strip() and not any(h in l.upper() for h in SECTION_HEADERS['summary'])]
            info['summary'] = ' '.join(body)[:700]
        if not info['summary']:
            for line in lines:
                if len(line) > 120:
                    info['summary'] = line[:500]
                    break

        # --- Education ---
        sec = get_section('education')
        degree_kw   = ['b.tech', 'btech', 'b.e', 'mtech', 'm.tech', 'bca', 'mca', 'mba',
                        'bsc', 'msc', 'bachelor', 'master', 'phd', 'computer science', 'engineering']
        inst_kw     = ['university', 'college', 'institute', 'school',
                        'manipal', 'iit', 'nit', 'bits', 'amity', 'vit']
        if sec:
            edu_lines = [l.strip() for l in sec.split('\n')
                         if l.strip() and not any(h in l.upper() for h in SECTION_HEADERS['education'])]
            deg, ins = '', ''
            for line in edu_lines[:20]:
                ll = line.lower()
                if any(k in ll for k in degree_kw) and not deg:
                    deg = line
                elif any(k in ll for k in inst_kw) and not ins:
                    ins = line
                if deg or ins:
                    if deg and ins:
                        info['education'].append((deg, ins))
                        deg, ins = '', ''
            if deg:
                info['education'].append((deg, ins))
        if not info['education']:
            info['education'] = [('Degree', 'University')]

        # --- Certifications ---
        sec = get_section('certifications')
        if sec:
            junk = set(h.upper() for hs in SECTION_HEADERS.values() for h in hs)
            for line in sec.split('\n'):
                cleaned = line.strip().lstrip('●•-– ').strip()
                upper_c = cleaned.upper()
                if (len(cleaned) > 8 and
                    upper_c not in junk and
                    upper_c != cleaned):   # skip ALL-CAPS header lines
                    info['certifications'].append(cleaned)

        # --- Skills ---
        sec = get_section('skills')
        if sec:
            collected = []
            for line in sec.split('\n'):
                line = line.strip()
                if not line or any(h in line.upper() for h in SECTION_HEADERS['skills']):
                    continue
                body = line.split(':', 1)[1] if ':' in line else line
                for tok in re.split(r'[,;]', body):
                    tok = tok.strip().lstrip('●•-– ').strip()
                    if 1 < len(tok) < 60:
                        collected.append(tok)
            info['skills'] = collected[:20]
        if not info['skills']:
            fallback_kw = ['Python', 'Java', 'SQL', 'AWS', 'Azure', 'Docker',
                           'Kubernetes', 'Git', 'MongoDB', 'React', 'FastAPI']
            found_skills = {k for k in fallback_kw if k.lower() in resume_text.lower()}
            info['skills'] = list(found_skills)[:12]

        # --- Work Experience ---
        sec = get_section('experience')
        if sec:
            role_kw = ['Engineer', 'Developer', 'Intern', 'Analyst', 'Architect',
                        'Manager', 'Lead', 'Consultant', 'Specialist', 'Designer']
            company_kw = ['Inc', 'Ltd', 'Pvt', 'Technologies', 'Accenture',
                           'Corp', 'Solutions', 'Services', 'Systems', 'Company']
            all_kw = role_kw + company_kw

            current_company = ''
            current_bullets: list = []

            for line in sec.split('\n'):
                stripped = line.strip()
                if not stripped:
                    continue
                is_header = any(h in stripped.upper() for h in SECTION_HEADERS['experience'])
                is_bullet = stripped[0] in ('●', '•', '-') if stripped else False
                is_company = (not is_bullet and not is_header and
                               5 < len(stripped) < 120 and
                               any(k in stripped for k in all_kw))

                if is_company:
                    if current_company and current_bullets:
                        info['experience_table'].append({
                            'company_name': current_company,
                            'roles_responsibility': current_bullets[:]
                        })
                    current_company = stripped
                    current_bullets = []
                elif is_bullet:
                    current_bullets.append(stripped.lstrip('●•-– ').strip())

            if current_company and current_bullets:
                info['experience_table'].append({
                    'company_name': current_company,
                    'roles_responsibility': current_bullets[:]
                })

        return info

# Professional Resume Formatter Class
class ProfessionalResumeFormatter:
    def __init__(self, logo_path=None):
        self.doc = Document()
        self.logo_path = logo_path
        self.blue_color = RGBColor(0, 51, 102)    # #003366 - for "Key Skills & Technologies" only
        self.black_color = RGBColor(0, 0, 0)      # #000000 - for all other text
        self.setup_document()
    
    def setup_document(self):
        """Set document margins"""
        for section in self.doc.sections:
            section.top_margin = Inches(0.5)
            section.bottom_margin = Inches(0.5)
            section.left_margin = Inches(0.75)
            section.right_margin = Inches(0.75)
    
    def add_header_with_logo(self, name, location):
        """Add header with logo and name (RIGHT aligned)"""
        header_table = self.doc.add_table(rows=1, cols=2)
        
        # Logo cell
        logo_cell = header_table.cell(0, 0)
        if self.logo_path and os.path.exists(self.logo_path):
            try:
                para = logo_cell.paragraphs[0]
                para.clear()
                run = para.add_run()
                run.add_picture(self.logo_path, width=Inches(2.5))
            except:
                logo_cell.text = "⊞ Logo"
        else:
            logo_cell.text = "⊞ Logo"
        
        # Name cell - RIGHT aligned (ONLY name, no location)
        name_cell = header_table.cell(0, 1)
        name_para = name_cell.paragraphs[0]
        name_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        name_run = name_para.add_run(name)
        set_calibri_font(name_run, 18, True, self.black_color)
        
        # Location removed from header - will be added separately in document body
        
        remove_table_borders(header_table)
        self.doc.add_paragraph()
    
    def add_page_header(self, name):
        """Add header for new pages (logo + name only)"""
        header_table = self.doc.add_table(rows=1, cols=2)
        
        # Logo cell
        logo_cell = header_table.cell(0, 0)
        if self.logo_path and os.path.exists(self.logo_path):
            try:
                para = logo_cell.paragraphs[0]
                para.clear()
                run = para.add_run()
                run.add_picture(self.logo_path, width=Inches(2.5))
            except:
                logo_cell.text = "⊞ Logo"
        else:
            logo_cell.text = "⊞ Logo"
        
        # Name cell - RIGHT aligned
        name_cell = header_table.cell(0, 1)
        name_para = name_cell.paragraphs[0]
        name_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        name_run = name_para.add_run(name)
        set_calibri_font(name_run, 18, True, self.black_color)
        
        remove_table_borders(header_table)
        self.doc.add_paragraph()
    
    def add_objective(self, text):
        """Add justified objective paragraph"""
        para = self.doc.add_paragraph()
        run = para.add_run(text)
        set_calibri_font(run, 11, False, self.black_color)
        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        para.space_after = Pt(4)  # Reduced from 12 to 6 points
    
    def add_section_header(self, text, use_blue=False):
        """Add section headers"""
        para = self.doc.add_paragraph()
        run = para.add_run(text)
        color = self.blue_color if use_blue else self.black_color
        set_calibri_font(run, 12, True, color)
        para.space_after = Pt(0)  # No spacing after section headers
    
    def add_blue_header(self, text):
        """Add blue background header like 'Key Expertise' and 'Education Details'"""
        # Create a table with 1 row and 1 column for the blue background
        header_table = self.doc.add_table(rows=1, cols=1)
        header_table.allow_autofit = False
        header_table.autofit = False
        
        # Set table width to use full page width
        header_table.columns[0].width = Inches(6.5)
        
        # Get the cell and set blue background
        cell = header_table.rows[0].cells[0]
        cell.text = text
        
        # Set darker light blue background
        shading = OxmlElement('w:shd')
        shading.set(qn('w:val'), 'clear')
        shading.set(qn('w:color'), 'auto')
        shading.set(qn('w:fill'), 'DCE6F0')  # Darker light blue color
        cell._element.get_or_add_tcPr().append(shading)
        
        # Format text (black, bold, left-aligned)
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            for run in paragraph.runs:
                set_calibri_font(run, 12, True, RGBColor(0, 0, 0))  # Black text
        
        # Remove table borders
        remove_table_borders(header_table)
        
        # Add minimal spacing after header
        self.doc.add_paragraph()
        # Reduce spacing after blue header
        last_para = self.doc.paragraphs[-1]
        last_para.space_after = Pt(3)  # Reduced spacing
    
    def add_page_header_to_all_pages(self, name):
        """Add header with logo and name that appears on every page"""
        # Get the section
        section = self.doc.sections[0]
        
        # Create header
        header = section.header
        header.paragraphs.clear()  # Clear any existing header content
        
        # Create header table with logo on left and name on right
        header_table = header.add_table(rows=1, cols=2, width=Inches(6.5))
        header_table.allow_autofit = False
        header_table.autofit = False
        
        # Set column widths - reduced gaps
        header_table.columns[0].width = Inches(3.0)  # Logo column - reduced from 3.5
        header_table.columns[1].width = Inches(3.5)  # Name column - increased from 3.0
        
        # Left cell - Logo
        logo_cell = header_table.rows[0].cells[0]
        logo_para = logo_cell.paragraphs[0]
        logo_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        
        # Add logo image if available
        if self.logo_path and os.path.exists(self.logo_path):
            try:
                logo_run = logo_para.add_run()
                logo_run.add_picture(self.logo_path, width=Inches(2.5))
            except Exception as e:
                print(f"Error adding logo to header: {e}")
                # Fallback to text if image fails
                logo_run = logo_para.add_run("⊞ Logo")
                set_calibri_font(logo_run, 10, False, self.black_color)
        else:
            # Fallback to text if no logo file
            logo_run = logo_para.add_run("⊞ Logo")
            set_calibri_font(logo_run, 10, False, self.black_color)
        
        # Right cell - Name
        name_cell = header_table.rows[0].cells[1]
        name_para = name_cell.paragraphs[0]
        name_run = name_para.add_run(name)
        set_calibri_font(name_run, 16, True, self.black_color)  # Increased from 12 to 16
        name_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        # Remove table borders
        remove_table_borders(header_table)
    
    def add_table(self, headers, data_rows):
        """Add table with labels on left, content on right"""
        try:
            print(f"Starting table creation...")
            print(f"Data rows: {data_rows}")
            print(f"Number of data rows: {len(data_rows)}")
            
            # Create table with 2 columns (labels on left, content on right)
            table = self.doc.add_table(rows=len(data_rows), cols=2)
            print(f"Table created successfully")
            
            # Set table style
            table.style = 'Table Grid'
            print(f"Table style set")
            
            # Set table properties
            table.autofit = False
            table.allow_autofit = False
            
            # Set column widths - labels on left, content on right
            table.columns[0].width = Inches(2.0)  # Labels column (Company Name, Roles & Responsibility)
            table.columns[1].width = Inches(5.0)  # Content column (actual company info and responsibilities)
            print(f"Column widths set")
            
            # Data rows - no header row needed
            for row_idx, row_data in enumerate(data_rows):
                print(f"Processing row {row_idx}: {row_data}")
                data_row = table.rows[row_idx]
                
                # Left cell - Label (Company Name, Roles & Responsibility)
                left_cell = data_row.cells[0]
                left_cell.text = row_data[0]  # Label text
                set_cell_background_white(left_cell)
                
                # Format left cell (labels) - bold text
                for paragraph in left_cell.paragraphs:
                    # Center-align "Roles & Responsibility", left-align "Company Name"
                    if row_data[0] == "Roles & Responsibility":
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        # Force center alignment by setting paragraph properties
                        paragraph.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    else:
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        paragraph.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in paragraph.runs:
                        set_calibri_font(run, 11, True, self.black_color)  # Bold for labels
                
                # Right cell - Content (company info or responsibilities)
                right_cell = data_row.cells[1]
                right_cell.text = ""
                
                # Handle multi-line text with bullet points
                if isinstance(row_data[1], str) and '\n' in row_data[1]:
                    lines = row_data[1].split('\n')
                    for i, line in enumerate(lines):
                        if line.strip():
                            if i == 0:
                                para = right_cell.paragraphs[0]
                                para.text = line.strip()
                            else:
                                para = right_cell.add_paragraph()
                                para.text = line.strip()
                else:
                    right_cell.text = str(row_data[1])
                
                # Format right cell (content)
                set_cell_background_white(right_cell)
                for paragraph in right_cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in paragraph.runs:
                        # Make company name bold, responsibilities regular
                        if row_data[0] == "Company Name":  # If this is the company name row
                            set_calibri_font(run, 11, True, self.black_color)  # Bold for company name
                        else:
                            set_calibri_font(run, 11, False, self.black_color)  # Regular text for responsibilities
                
                print(f"Row {row_idx} processed")
            
            # Add spacing after table
            self.doc.add_paragraph()
            print(f"Table creation completed successfully")
            
        except Exception as e:
            print(f"Error creating table: {str(e)}")
            import traceback
            traceback.print_exc()
            # Fallback: add simple text instead of table
            self.doc.add_paragraph("Error creating table - displaying as text:")
            for row_data in data_rows:
                self.doc.add_paragraph(f"{row_data[0]}: {row_data[1]}")
                self.doc.add_paragraph()
    
    def add_standard_table(self, headers, data_rows):
        """Add standard table with headers at top (for education, etc.)"""
        try:
            print(f"Starting standard table creation...")
            print(f"Headers: {headers}")
            print(f"Data rows: {data_rows}")
            print(f"Number of data rows: {len(data_rows)}")
            
            # Create table with header row + data rows
            table = self.doc.add_table(rows=len(data_rows) + 1, cols=len(headers))
            print(f"Table created successfully")
            
            # Set table style
            table.style = 'Table Grid'
            print(f"Table style set")
            
            # Set table properties
            table.autofit = False
            table.allow_autofit = False
            
            # Set column widths
            table.columns[0].width = Inches(3.0)  # Qualification column
            table.columns[1].width = Inches(4.0)  # Institution column
            print(f"Column widths set")
            
            # Header row
            header_row = table.rows[0]
            for i, header in enumerate(headers):
                cell = header_row.cells[i]
                cell.text = header
                set_cell_background_white(cell)
                
                # Format header text - bold
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in paragraph.runs:
                        set_calibri_font(run, 11, True, self.black_color)
            print(f"Header row created")
            
            # Data rows
            for row_idx, row_data in enumerate(data_rows, 1):
                print(f"Processing row {row_idx}: {row_data}")
                data_row = table.rows[row_idx]
                for col_idx, cell_text in enumerate(row_data):
                    cell = data_row.cells[col_idx]
                    cell.text = str(cell_text)
                    
                    # Format cell
                    set_cell_background_white(cell)
                    for paragraph in cell.paragraphs:
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        for run in paragraph.runs:
                            set_calibri_font(run, 11, False, self.black_color)
                print(f"Row {row_idx} processed")
            
            # Add spacing after table
            self.doc.add_paragraph()
            print(f"Standard table creation completed successfully")
            
        except Exception as e:
            print(f"Error creating standard table: {str(e)}")
            import traceback
            traceback.print_exc()
            # Fallback: add simple text instead of table
            self.doc.add_paragraph("Error creating standard table - displaying as text:")
            for row_data in data_rows:
                self.doc.add_paragraph(f"{headers[0]}: {row_data[0]}, {headers[1]}: {row_data[1]}")
                self.doc.add_paragraph()
    
    def add_bullet_list(self, items):
        """Add bullet point list with minimal spacing"""
        for item in items:
            para = self.doc.add_paragraph()
            run = para.add_run(f"• {item}")
            set_calibri_font(run, 11, False, self.black_color)
            # Eliminate all spacing between bullet points
            para.space_after = Pt(0)  # No spacing after paragraph
            para.space_before = Pt(0)  # No spacing before paragraph
            # Set line spacing to single
            para.paragraph_format.line_spacing = 1.0
    
    def create_complete_resume(self, candidate_info):
        """Create complete professional resume"""
        # Add page header that appears on every page
        self.add_page_header_to_all_pages(candidate_info['name'])
        
        # Objective/Summary (only if present)
        if candidate_info.get('summary'):
            self.add_blue_header("Key Expertise")
            self.add_objective(candidate_info['summary'])
        
        # Education table (only if present)
        if candidate_info.get('education'):
            self.add_blue_header("Education Details")
            self.add_standard_table(["Qualification", "Institution"], candidate_info['education'])
        
        # Contact info - compact spacing
        name_para = self.doc.add_paragraph()
        name_run = name_para.add_run(candidate_info['name'])
        set_calibri_font(name_run, 16, True, self.black_color)
        
        if candidate_info.get('location'):
            loc_para = self.doc.add_paragraph()
            loc_run = loc_para.add_run(candidate_info['location'])
            set_calibri_font(loc_run, 11, False, self.black_color)
            loc_para.space_after = Pt(3)  # Minimal spacing after location
        
        # Date and Subject - compact spacing
        if candidate_info.get('date'):
            date_para = self.doc.add_paragraph()
            date_run = date_para.add_run(f"Date: {candidate_info['date']}")
            set_calibri_font(date_run, 11, False, self.black_color)
            date_para.space_after = Pt(3)  # Minimal spacing after date
        
        if candidate_info.get('subject'):
            subj_para = self.doc.add_paragraph()
            subj_label = subj_para.add_run("Subject: ")
            set_calibri_font(subj_label, 11, True, self.black_color)
            subj_text = subj_para.add_run(str(candidate_info['subject']))
            set_calibri_font(subj_text, 11, False, self.black_color)
            subj_para.space_after = Pt(3)  # Minimal spacing after subject
        
        # Cover letter (only if present)
        if candidate_info.get('cover_letter'):
            dear_para = self.doc.add_paragraph()
            dear_run = dear_para.add_run("Dear Hiring Manager,")
            set_calibri_font(dear_run, 11, True, self.black_color)  # Made bold
            # Split cover letter into paragraphs and add each separately
            cover_paragraphs = candidate_info['cover_letter'].split('\n\n')
            for para_text in cover_paragraphs:
                if para_text.strip() and "Dear Hiring Manager" not in para_text.strip():
                    cover_para = self.doc.add_paragraph()
                    cover_run = cover_para.add_run(para_text.strip())
                    set_calibri_font(cover_run, 11, False, self.black_color)
                    cover_para.space_after = Pt(3)  # Reduced spacing
        
        # Skills section - Ultra compact spacing like certifications
        if candidate_info.get('skills'):
            self.add_section_header("Key Skills")
            # Add skills with ultra-compact spacing
            for item in candidate_info['skills']:
                para = self.doc.add_paragraph()
                run = para.add_run(f"\t• {item}")
                set_calibri_font(run, 11, False, self.black_color)
                # Ultra-compact spacing for skills
                para.space_after = Pt(0)  # No spacing after
                para.space_before = Pt(0)  # No spacing before
                para.paragraph_format.line_spacing = 1.0  # Single line spacing
                para.paragraph_format.space_after = Pt(0)  # Additional Word-level spacing control
            # Insert a single blank line after skills list
            self.doc.add_paragraph()
 
        # Experience section
        if candidate_info.get('experience_table'):
            self.add_section_header("Experience:")
            print(f"Found {len(candidate_info['experience_table'])} experience entries")
            print(f"Experience table data: {candidate_info['experience_table']}")
            
            for exp in candidate_info['experience_table']:
                print(f"Processing experience entry: {exp}")
                
                # Convert array of responsibilities to bullet points
                responsibilities_text = ""
                if isinstance(exp['roles_responsibility'], list):
                    for i, resp in enumerate(exp['roles_responsibility'], 1):
                        responsibilities_text += f"• {resp}\n"
                    responsibilities_text = responsibilities_text.strip()
                else:
                    responsibilities_text = str(exp['roles_responsibility'])
                
                print(f"Company: {exp['company_name']}")
                print(f"Responsibilities: {responsibilities_text[:100]}...")
                
                # Create table with labels on left, content on right
                exp_data = [
                    ("Company Name", exp['company_name']),
                    ("Roles & Responsibility", responsibilities_text)
                ]
                print(f"Creating table with data: {exp_data}")
                self.add_table(["", ""], exp_data)  # Empty headers since labels are in left column
                # Spacing after table is already added inside add_table(); avoid extra spacing
        elif candidate_info.get('experience'):  # Fallback for old format
            self.add_section_header("Experience:")
            print(f"Using fallback experience format: {candidate_info['experience']}")
            for exp in candidate_info['experience'][:3]:
                # Convert array of responsibilities to bullet points
                responsibilities_text = ""
                if isinstance(exp['responsibilities'], list):
                    for i, resp in enumerate(exp['responsibilities'], 1):
                        responsibilities_text += f"• {resp}\n"
                    responsibilities_text = responsibilities_text.strip()
                else:
                    responsibilities_text = str(exp['responsibilities'])
                
                # Create table with labels on left, content on right
                exp_data = [
                    ("Company Name", f"{exp['company']} as {exp['role']} ({exp['duration']})"),
                    ("Roles & Responsibility", responsibilities_text)
                ]
                self.add_table(["", ""], exp_data)  # Empty headers since labels are in left column
                # Spacing after table is already added inside add_table(); avoid extra spacing
        
        # Certifications - ultra compact spacing
        if candidate_info.get('certifications'):
            self.add_section_header("CERTIFICATIONS")
            # Add certifications with ultra-compact spacing
            for item in candidate_info['certifications']:
                para = self.doc.add_paragraph()
                run = para.add_run(f"• {item}")
                set_calibri_font(run, 11, False, self.black_color)
                # Ultra-compact spacing for certifications
                para.space_after = Pt(0)  # No spacing after
                para.space_before = Pt(0)  # No spacing before
                para.paragraph_format.line_spacing = 1.0  # Single line spacing
                para.paragraph_format.space_after = Pt(0)  # Additional spacing control
        
        # Debug: Check if we have any experience data at all
        if not candidate_info.get('experience_table') and not candidate_info.get('experience'):
            print("WARNING: No experience data found at all!")
            print(f"All available keys: {list(candidate_info.keys())}")
            self.doc.add_paragraph("No experience data available")
        
        return self.save_to_buffer()
    
    def save_to_buffer(self):
        """Save to memory buffer"""
        buffer = io.BytesIO()
        self.doc.save(buffer)
        buffer.seek(0)
        return buffer

# Professional PPTX Formatter Class
class ProfessionalPPTXFormatter:
    def __init__(self, logo_path=None):
        self.prs = Presentation()
        # Blank slide layout is index 6 usually
        self.slide_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(self.slide_layout)
        # Using 13.33x7.5 for widescreen
        self.prs.slide_width = PptxInches(13.33)
        self.prs.slide_height = PptxInches(7.5)
        self.logo_path = logo_path
        self.green_color = PptxRGBColor(50, 150, 50) # Soft green
        self.blue_color = PptxRGBColor(100, 150, 200) # Soft blue
        
    def add_header(self, name, role, domain):
        txBox = self.slide.shapes.add_textbox(PptxInches(0.4), PptxInches(0.3), PptxInches(12.5), PptxInches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        if domain:
            text = f"{name} - {role[:50]} (specialized in {domain[:30]})"
        else:
            text = f"{name} - {role[:50]}"
        p.text = text
        p.font.size = PptxPt(24)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0,0,0)
        
    def add_left_column(self, candidate_info):
        left_margin = PptxInches(0.4)
        box_width = PptxInches(4.2)
        
        # 1. Contact Box
        shape_contact = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_margin, PptxInches(1.2), box_width, PptxInches(1.5)
        )
        shape_contact.fill.background()
        shape_contact.line.color.rgb = self.blue_color
        shape_contact.line.width = PptxPt(1.5)
        
        tf = shape_contact.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        # Small top margin so text doesn't touch the border
        tf.margin_top = PptxPt(8)
        tf.margin_left = PptxPt(8)
        tf.margin_bottom = PptxPt(4)
        
        info_lines = [
            f"Name: {candidate_info.get('name', '')}",
            f"Role: {candidate_info.get('subject', '').replace('Application for the Position of', '').strip()[:50]}",
            f"Email: {candidate_info.get('email', '')}"
        ]
        if candidate_info.get('phone'):
             info_lines.append(f"Phone: {candidate_info.get('phone', '')}")
             
        is_first = True
        for line in info_lines:
            if line.replace('Phone: ', '').replace('Email: ', '').replace('Role: ', '').strip():
                 if is_first:
                     p = tf.paragraphs[0]
                     is_first = False
                 else:
                     p = tf.add_paragraph()
                 p.text = line
                 p.font.size = PptxPt(11)
                 p.font.color.rgb = PptxRGBColor(0,0,0)
        
        # 2. Professional Summary
        shape_summary = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_margin, PptxInches(2.9), box_width, PptxInches(2.4)
        )
        shape_summary.fill.background()
        shape_summary.line.color.rgb = self.green_color
        shape_summary.line.width = PptxPt(1.5)
        
        txBox_title1 = self.slide.shapes.add_textbox(left_margin, PptxInches(2.9), box_width, PptxInches(0.4))
        p_t1 = txBox_title1.text_frame.paragraphs[0]
        p_t1.text = "Professional Summary"
        p_t1.font.bold = True
        p_t1.font.size = PptxPt(12)
        p_t1.alignment = PP_ALIGN.CENTER
        
        txBox_body1 = self.slide.shapes.add_textbox(left_margin+PptxInches(0.1), PptxInches(3.3), box_width-PptxInches(0.2), PptxInches(1.9))
        txBox_body1.text_frame.word_wrap = True
        p_b1 = txBox_body1.text_frame.paragraphs[0]
        p_b1.text = candidate_info.get('summary', '')
        p_b1.font.size = PptxPt(10)
        p_b1.font.color.rgb = PptxRGBColor(0,0,0)
        
        # 3. Education / Certification
        txBox_edu = self.slide.shapes.add_textbox(left_margin, PptxInches(5.4), box_width, PptxInches(0.4))
        p_t2 = txBox_edu.text_frame.paragraphs[0]
        p_t2.text = "Education / Certification"
        p_t2.font.bold = True
        p_t2.font.size = PptxPt(12)
        p_t2.alignment = PP_ALIGN.CENTER
        
        txBox_body2 = self.slide.shapes.add_textbox(left_margin, PptxInches(5.8), box_width, PptxInches(1.5))
        txBox_body2.text_frame.word_wrap = True
        
        is_first = True
        edu_list = candidate_info.get('education', [])
        for edu in edu_list[:2]:
            if is_first:
                p_e = txBox_body2.text_frame.paragraphs[0]
                is_first = False
            else:
                p_e = txBox_body2.text_frame.add_paragraph()
            p_e.text = f"• {edu[0]} - {edu[1]}"
            p_e.font.size = PptxPt(10)
            
        cert_list = candidate_info.get('certifications', [])
        for cert in cert_list[:3]:
            if is_first:
                p_c = txBox_body2.text_frame.paragraphs[0]
                is_first = False
            else:
                p_c = txBox_body2.text_frame.add_paragraph()
            p_c.text = f"• {cert}"
            p_c.font.size = PptxPt(10)

    def add_right_column(self, candidate_info):
        right_margin = PptxInches(4.8)
        box_width = PptxInches(8.1)
        
        # 1. Core Skills Table
        txBox_skills = self.slide.shapes.add_textbox(right_margin, PptxInches(1.15), box_width, PptxInches(0.35))
        p_sk = txBox_skills.text_frame.paragraphs[0]
        p_sk.text = "Core Skills"
        p_sk.font.bold = True
        p_sk.font.size = PptxPt(12)
        p_sk.alignment = PP_ALIGN.CENTER
        
        table_shape = self.slide.shapes.add_table(4, 2, right_margin, PptxInches(1.5), box_width, PptxInches(1.5))
        table = table_shape.table
        table.columns[0].width = PptxInches(2.5)
        table.columns[1].width = PptxInches(5.6)
        
        # Distribute skills across 4 rows
        all_skills = candidate_info.get('skills', [])
        n = len(all_skills)
        # Split into 4 roughly equal quarters
        q = max(1, n // 4)
        tech_skills       = ", ".join(all_skills[:q])
        frameworks        = ", ".join(all_skills[q:q*2])
        tools             = ", ".join(all_skills[q*2:q*3])
        domain_exp        = ", ".join(all_skills[q*3:])

        rows_data = [
            ("Technical Skills",      tech_skills),
            ("Frameworks/Libraries",  frameworks),
            ("Tools",                 tools),
            ("Domain Experience",     domain_exp),
        ]
        
        for i, (cat, val) in enumerate(rows_data):
            cell_left = table.cell(i, 0)
            cell_left.text = cat
            cell_left.text_frame.paragraphs[0].font.bold = True
            cell_left.text_frame.paragraphs[0].font.size = PptxPt(10)
            cell_left.text_frame.word_wrap = True

            cell_right = table.cell(i, 1)
            cell_right.text = val
            cell_right.text_frame.paragraphs[0].font.size = PptxPt(9)
            cell_right.text_frame.word_wrap = True

        # Fixed y for Responsibilities — always below the table
        resp_label_y = PptxInches(3.3)
        resp_box_y   = PptxInches(3.65)

        # 2. Responsibilities label + rounded rectangle
        txBox_resp = self.slide.shapes.add_textbox(right_margin, resp_label_y, box_width, PptxInches(0.35))
        p_r = txBox_resp.text_frame.paragraphs[0]
        p_r.text = "Responsibilities :"
        p_r.font.bold = True
        p_r.font.size = PptxPt(12)

        shape_resp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, right_margin, resp_box_y, box_width, PptxInches(3.5)
        )
        shape_resp.fill.background()
        shape_resp.line.color.rgb = self.green_color
        shape_resp.line.width = PptxPt(1.5)

        # Use a separate textbox OVER the shape so text stays inside the box bounds
        resp_inner = self.slide.shapes.add_textbox(
            right_margin + PptxInches(0.12),
            resp_box_y + PptxInches(0.1),
            box_width - PptxInches(0.24),
            PptxInches(3.3)
        )
        tf_r = resp_inner.text_frame
        tf_r.word_wrap = True

        exp_list = candidate_info.get('experience_table', [])
        if not exp_list and candidate_info.get('experience'):
            exp_list = candidate_info.get('experience')

        is_first = True

        for exp in exp_list:
            resp = exp.get('roles_responsibility') or exp.get('responsibilities')
            items = resp if isinstance(resp, list) else [l.strip() for l in (resp or '').split('\n') if l.strip()]
            for r in items:
                _r_clean = r.lstrip('•– ').strip()
                text = f"• {_r_clean}"
                if is_first:
                    p = tf_r.paragraphs[0]
                    is_first = False
                else:
                    p = tf_r.add_paragraph()
                p.text = text
                p.font.size = PptxPt(10)
                p.font.color.rgb = PptxRGBColor(0, 0, 0)
                # small space between bullets
                p.space_after = PptxPt(2)

    def create_complete_resume(self, candidate_info):
        role = candidate_info.get('subject', '').replace('Application for the Position of', '').strip()
        skills = candidate_info.get('skills', [])

        # Extract a meaningful domain from role — strip generic job title words
        generic_words = {'Developer', 'Engineer', 'Manager', 'Analyst', 'Lead',
                         'Senior', 'Junior', 'Associate', 'Consultant', 'Specialist',
                         'Architect', 'Designer', 'Administrator', 'Officer', 'Executive'}
        domain_words = [w for w in role.split() if w not in generic_words]
        if domain_words:
            domain = ' '.join(domain_words[:3])  # e.g. "ServiceNow" or "AI/ML"
        elif skills:
            domain = ', '.join(skills[:2])
        else:
            domain = 'Technology'

        self.add_header(candidate_info.get('name', 'Name'), role, domain)
        self.add_left_column(candidate_info)
        self.add_right_column(candidate_info)
        return self.save_to_buffer()

    def save_to_buffer(self):
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

# Text Extraction Classes
class ResumeTextExtractor:
    def extract_from_pdf(self, file_content):
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text is None:
                    page_text = ""
                text += page_text + "\n"
            return text
        except Exception as e:
            return f"PDF Error: {str(e)}"
    
    def extract_from_docx(self, file_content):
        try:
            doc = Document(io.BytesIO(file_content))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
                
            # ✅ NEW: Extract text from tables also
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        text += " | ".join(row_data) + "\n"
                        
            return text
        except Exception as e:
            return f"DOCX Error: {str(e)}"
    
    def extract_from_txt(self, file_content):
        try:
            return file_content.decode('utf-8', errors='ignore')
        except Exception as e:
            return f"TXT Error: {str(e)}"

def extract_text_from_file(file):
    """Extract text from uploaded file"""
    extractor = ResumeTextExtractor()
    content = file.read()
    filename = file.filename.lower()
    
    if filename.endswith('.pdf'):
        return extractor.extract_from_pdf(content)
    elif filename.endswith(('.doc', '.docx')):
        return extractor.extract_from_docx(content)
    elif filename.endswith('.txt'):
        return extractor.extract_from_txt(content)
    else:
        return "Unsupported file format"

PROGRESS_PAGE = '''
<!DOCTYPE html>
<html>
<head>
    <title>Processing Resumes...</title>
    <meta charset="utf-8">
    <style>
        body { font-family: Arial, sans-serif; display: flex; justify-content: center; align-items: center; min-height: 100vh; margin: 0; background: #f5f5f5; }
        .card { background: #fff; border-radius: 12px; padding: 40px 50px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); text-align: center; width: 420px; }
        h2 { margin-bottom: 8px; color: #1e3a6e; }
        .sub { color: #666; margin-bottom: 30px; font-size: 14px; }
        .bar-wrap { background: #e0e0e0; border-radius: 20px; height: 18px; overflow: hidden; margin-bottom: 14px; }
        .bar { height: 100%; background: linear-gradient(90deg, #1e3a6e, #3278c8); border-radius: 20px; transition: width 0.4s ease; width: 0%; }
        .count { font-size: 15px; color: #333; margin-bottom: 6px; }
        .errors { margin-top: 16px; text-align: left; font-size: 13px; color: #c0392b; }
        .done-btn { display: none; margin-top: 24px; background: #1e3a6e; color: #fff; border: none; border-radius: 8px; padding: 12px 32px; font-size: 16px; cursor: pointer; text-decoration: none; }
        .done-btn:hover { background: #3278c8; }
    </style>
</head>
<body>
<div class="card">
    <h2>Processing Resumes</h2>
    <p class="sub">Each resume is processed independently via a separate API call.</p>
    <div class="bar-wrap"><div class="bar" id="bar"></div></div>
    <div class="count" id="count">0 / {{ total }} done</div>
    <div class="errors" id="errors"></div>
    <a class="done-btn" id="dlBtn" href="/download/{{ job_id }}">Download ZIP</a>
</div>
<script>
    const jobId = "{{ job_id }}";
    const total = {{ total }};
    function poll() {
        fetch('/status/' + jobId)
            .then(r => r.json())
            .then(data => {
                const pct = total > 0 ? (data.done / total * 100) : 0;
                document.getElementById('bar').style.width = pct + '%';
                document.getElementById('count').textContent = data.done + ' / ' + total + ' done';
                if (data.errors && data.errors.length) {
                    const errDiv = document.getElementById('errors');
                    errDiv.textContent = '';
                    const label = document.createElement('b');
                    label.textContent = 'Errors:';
                    errDiv.appendChild(label);
                    data.errors.forEach(function(e) {
                        const d = document.createElement('div');
                        d.textContent = e;
                        errDiv.appendChild(d);
                    });
                }
                if (data.finished) {
                    document.getElementById('dlBtn').style.display = 'inline-block';
                } else {
                    setTimeout(poll, 1500);
                }
            })
            .catch(() => setTimeout(poll, 2000));
    }
    poll();
</script>
</body>
</html>
'''

# Main Flask Route
@app.route('/', methods=['GET', 'POST'])
def upload_and_format():
    if request.method == 'POST':
        resume_files = request.files.getlist('resume_file')
        resume_files = [f for f in resume_files if f.filename != '']
        if not resume_files:
            flash('Please select at least one resume file')
            return render_template_string(UPLOAD_FORM)

        logo_path = None

        def process_single(resume_file):
            """Process one resume file and return (buffer, safe_filename) or raise."""
            resume_text = extract_text_from_file(resume_file)
            print("\n" + "="*80)
            print(f"RAW EXTRACTED TEXT: {resume_file.filename}")
            print("="*80)
            print(resume_text[:1000])
            print("="*80 + "\n")

            if resume_text.startswith(('PDF Error', 'DOCX Error', 'TXT Error', 'Unsupported')):
                raise ValueError(f'Error processing {resume_file.filename}: {resume_text}')
            if len(resume_text.strip()) < 50:
                raise ValueError(f'{resume_file.filename} appears to be empty or too short')

            api_key = os.getenv('OPENAI_API_KEY')
            extractor = OpenAIResumeExtractor(api_key)
            candidate_info = extractor.extract_with_openai(resume_text)

            formatter = ProfessionalPPTXFormatter(logo_path)
            buffer = formatter.create_complete_resume(candidate_info)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            safe_name = re.sub(r'[^\w\s-]', '', candidate_info.get('name', 'resume')).strip().replace(' ', '_')
            filename = f"{safe_name}_{timestamp}.pptx"
            return buffer, filename

        try:
            if len(resume_files) == 1:
                # Single file — return PPTX directly
                buffer, filename = process_single(resume_files[0])
                return send_file(
                    buffer,
                    as_attachment=True,
                    download_name=filename,
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            else:
                # Multiple files — process in parallel threads, return job ID immediately
                job_id = str(uuid.uuid4())
                # Read file bytes now (can't pass Werkzeug file objects across threads)
                file_items = [(f.filename, f.read()) for f in resume_files]
                total = len(file_items)

                with jobs_lock:
                    jobs[job_id] = {
                        'total': total,
                        'done': 0,
                        'errors': [],
                        'zip_buffer': None,
                        'zip_name': None,
                        'lock': threading.Lock(),
                        'finished': False,
                    }

                def run_bulk(job_id, file_items, logo_path):
                    job = jobs[job_id]
                    zip_buffer = io.BytesIO()
                    zip_name = f"formatted_resumes_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"

                    def process_item(item):
                        orig_name, file_bytes = item
                        fake_file = io.BytesIO(file_bytes)
                        fake_file.filename = orig_name
                        return process_single(fake_file)

                    results = {}
                    with ThreadPoolExecutor(max_workers=4) as executor:
                        future_to_name = {executor.submit(process_item, item): item[0] for item in file_items}
                        for future in as_completed(future_to_name):
                            orig_name = future_to_name[future]
                            with job['lock']:
                                try:
                                    pptx_buffer, pptx_name = future.result()
                                    results[pptx_name] = pptx_buffer.read()
                                except Exception as err:
                                    job['errors'].append(f"{orig_name}: {str(err)}")
                                job['done'] += 1

                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for name, data in results.items():
                            zf.writestr(name, data)
                    zip_buffer.seek(0)

                    with job['lock']:
                        job['zip_buffer'] = zip_buffer
                        job['zip_name'] = zip_name
                        job['finished'] = True

                threading.Thread(target=run_bulk, args=(job_id, file_items, logo_path), daemon=True).start()

                return render_template_string(PROGRESS_PAGE, job_id=job_id, total=total)

        except Exception as e:
            flash(f'Processing error: {str(e)}')
            return render_template_string(UPLOAD_FORM)

    return render_template_string(UPLOAD_FORM)


@app.route('/status/<job_id>')
def job_status(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job not found'}), 404
    with job['lock']:
        return jsonify({
            'total': job['total'],
            'done': job['done'],
            'errors': job['errors'],
            'finished': job['finished'],
        })


@app.route('/download/<job_id>')
def job_download(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return 'Job not found', 404
    with job['lock']:
        if not job['finished']:
            return 'Still processing', 202
        zip_buffer = job['zip_buffer']
        zip_name = job['zip_name']
    # Clean up after download
    with jobs_lock:
        jobs.pop(job_id, None)
    zip_buffer.seek(0)
    return send_file(
        zip_buffer,
        as_attachment=True,
        download_name=zip_name,
        mimetype='application/zip'
    )


if __name__ == '__main__':
    import socket
    
    # Get local IP address for network access
    def get_local_ip():
        try:
            # Connect to a remote address to get local IP
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]
            s.close()
            return local_ip
        except:
            return "127.0.0.1"
    
    local_ip = get_local_ip()
    port = 5005
    
    print("=" * 60)
    print("🌐 RESUME FORMATTER - NETWORK ACCESS")
    print("=" * 60)
    print(f"Local access: http://localhost:{port}")
    print(f"Network access: http://{local_ip}:{port}")
    print(f"Other users on same network can use: http://{local_ip}:{port}")
    print("=" * 60)
    print("Starting server...")

    # Start ngrok tunnel if pyngrok is installed
    if ngrok is not None:
        try:
            ngrok_authtoken = os.getenv('NGROK_AUTHTOKEN')
            if ngrok_authtoken:
                ngrok.set_auth_token(ngrok_authtoken)
            http_tunnel = ngrok.connect(addr=port, proto="http")
            public_url = http_tunnel.public_url
            print("=" * 60)
            print(f"Public (ngrok) URL: {public_url}")
            print("Share this URL to access your app over the internet.")
            print("=" * 60)
        except Exception as e:
            print(f"ngrok start failed: {e}")
    else:
        print("pyngrok not installed; you can run 'ngrok http {port}' separately to expose the app.")
    
    port = int(os.environ.get('PORT', port))
    debug = os.environ.get('FLASK_ENV') != 'production'
    app.run(debug=debug, host='0.0.0.0', port=port, use_reloader=False)
